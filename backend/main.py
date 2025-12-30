from __future__ import annotations
# ===== OCR SUPPORT =====
import pytesseract
from PIL import Image
from pdf2image import convert_from_path

def extract_text_from_image_file(file_path: Path) -> str:
    """Ekstrak teks dari file gambar (jpg/png)."""
    try:
        img = Image.open(file_path)
        text = pytesseract.image_to_string(img, lang='ind')
        return text
    except Exception as e:
        print(f"[ERROR] OCR image: {e}")
        return ""

def extract_text_from_pdf_with_ocr(file_path: Path) -> str:
    """Ekstrak teks dari PDF (OCR jika p+erlu)."""
    try:
        pages = convert_from_path(str(file_path))
        text = ""
        for page in pages:
            text += pytesseract.image_to_string(page, lang='ind') + "\n"
        return text
    except Exception as e:
        print(f"[ERROR] OCR PDF: {e}")
        return ""

# ===== stdlib =====
from pathlib import Path
from datetime import datetime
from typing import List, Dict, Optional
import csv, json, subprocess, sys, random, re
import uuid # Added for UUID generation

# ===== third-party =====
from fastapi import FastAPI, HTTPException, Depends, UploadFile, File, Form
from fastapi.responses import FileResponse
from fastapi.middleware.cors import CORSMiddleware
from fastapi.staticfiles import StaticFiles
from fastapi.security import HTTPBearer, HTTPAuthorizationCredentials
from pydantic import BaseModel, Field, EmailStr
import joblib  # untuk load model .joblib
from passlib.context import CryptContext
from sqlalchemy.orm import Session # Added for DB session
from database import engine, get_db # Added for DB access
import models # Added for DB models
import schemas # Added for Pydantic schemas

import secrets
import os
import shutil

# ===== ENV VARS =====

# ===== ENV VARS =====
GOOGLE_GEMINI_API_KEY = os.getenv("GOOGLE_GEMINI_API_KEY")

# ===== Google Gemini AI =====
from google import genai


# ===== Google Gemini AI configure =====
CLIENT = genai.Client(api_key=GOOGLE_GEMINI_API_KEY) if GOOGLE_GEMINI_API_KEY else None

def _pick_available_model(preferred: list[str] | None = None) -> str:
    preferred = preferred or [
        "gemini-1.5-flash",
        "gemini-1.5-flash-latest",
        "gemini-1.5-flash-8b",
        "gemini-1.5-pro",
    ]
    try:
        models = CLIENT.models.list()
        names = [m.name for m in models] if models else []
        for p in preferred:
            if p in names:
                return p
        for n in names:
            if "gemini" in n and "flash" in n:
                return n
    except Exception as e:
        print(f"[AI] List models failed: {e}")
    return "gemini-1.5-flash"

# ===== pastikan package "training" bisa di-import =====
ROOT = Path(__file__).resolve().parents[1]  # .../LearnCheck
if str(ROOT) not in sys.path:
    sys.path.insert(0, str(ROOT))

# util LearnCheck
from training.scripts import common as LC

# ===== path & konstanta =====
SOAL_DIR = LC.SOAL_DIR
TRAINING_DIR = LC.TRAINING_DIR
JAWABAN_DIR = LC.JAWABAN_DIR
MODELS_DIR = LC.MODELS_OUTPUT_DIR            # tempat file model *.joblib
MATERI_DIR = TRAINING_DIR / "materi"
UPLOADS_DIR = Path(__file__).resolve().parents[1] / "uploads"
UPLOADS_DIR.mkdir(exist_ok=True)
CHOICES = {"A", "B", "C", "D", "E"}

# ===== FastAPI app =====
app = FastAPI(title="LearnCheck API", version="0.4.0")
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],   # batasi ke domain FE saat deploy
    allow_methods=["*"],
    allow_headers=["*"],
)

# Serve static files for uploaded materials
# app.mount("/uploads", StaticFiles(directory=str(UPLOADS_DIR)), name="uploads")

@app.get("/uploads/{filename}")
async def get_uploaded_file(filename: str):
    """Explicit endpoint to serve uploaded files."""
    file_path = UPLOADS_DIR / filename
    print(f"[SERVE] Request for file: {filename}")
    print(f"[SERVE] Checking path: {file_path}")
    if not file_path.exists():
        print(f"[SERVE] File not found: {file_path}")
        raise HTTPException(status_code=404, detail="File not found on server")
    return FileResponse(file_path)

# ===== Database imports =====
from backend.db import engine, Base, get_db
from backend.models import User as DBUser, Class as DBClass, Session as DBSession, Quiz as DBQuiz, Material as DBMaterial
import backend.crud as crud
import backend.schemas as schemas
from sqlalchemy.orm import Session as DBSessionType
from sqlalchemy import text
import json

# Create tables on startup
@app.on_event("startup")
async def startup_event():
    Base.metadata.create_all(bind=engine)
    # Ensure missing columns exist (for runtime migrations)
    try:
        from sqlalchemy import text as _text
        with engine.begin() as conn:
            conn.execute(_text("ALTER TABLE materials ADD COLUMN IF NOT EXISTS file_type VARCHAR(50)"))
            conn.execute(_text("ALTER TABLE materials ADD COLUMN IF NOT EXISTS created_at TIMESTAMP DEFAULT NOW()"))
            conn.execute(_text("ALTER TABLE materials ADD COLUMN IF NOT EXISTS uploader_id INTEGER"))
            conn.execute(_text("ALTER TABLE materials ADD COLUMN IF NOT EXISTS file_url VARCHAR(255)"))
            # Migrate legacy column 'uploaded_by' to new 'uploader_id' if exists
            has_legacy = conn.execute(
                _text("SELECT 1 FROM information_schema.columns WHERE table_name='materials' AND column_name='uploaded_by'")
            ).fetchone()
            if has_legacy:
                # Drop NOT NULL to avoid violation during migration
                try:
                    conn.execute(_text("ALTER TABLE materials ALTER COLUMN uploaded_by DROP NOT NULL"))
                except Exception:
                    pass
                conn.execute(_text("UPDATE materials SET uploader_id = uploaded_by WHERE uploader_id IS NULL"))
            # Rename legacy column 'durl' to 'file_url' if present
            has_durl = conn.execute(
                _text("SELECT 1 FROM information_schema.columns WHERE table_name='materials' AND column_name='durl'")
            ).fetchone()
            if has_durl:
                try:
                    conn.execute(_text("ALTER TABLE materials RENAME COLUMN durl TO file_url"))
                except Exception:
                    pass
    except Exception as e:
        print(f"[STARTUP] Warn: failed to ensure materials columns: {e}")
    # Create default admin if not exists
    db = next(get_db())
    admin = db.query(DBUser).filter(DBUser.email == "admin@learncheck.com").first()
    if not admin:
        admin = DBUser(
            email="admin@learncheck.com",
            username="Admin",
            hashed_password=hash_password("admin123"),
            role="admin"
        )
        db.add(admin)
        db.commit()
        print("[STARTUP] Default admin created")
    db.close()

# ===== cache di memori =====
BANKS: Dict[str, dict] = {}
ANSKEY: Dict[tuple, dict] = {}   # (mapel, qid) -> {kunci,bobot,topik,tingkat}
ID2MAPEL: Dict[str, set] = {}    # qid -> {mapel}
MODELS: Dict[str, object] = {}   # mapel -> sklearn pipeline

# ===== User Management =====
pwd_context = CryptContext(schemes=["bcrypt"], deprecated="auto")
security = HTTPBearer()

def hash_password(password: str) -> str:
    return pwd_context.hash(password)

def verify_password(plain: str, hashed: str) -> bool:
    return pwd_context.verify(plain, hashed)

def create_token() -> str:
    return secrets.token_urlsafe(32)

# ===== helpers (token & normalisasi) =====
STOP = set("yang dari di ke dan untuk pada adalah ini itu tersebut dengan sebagai tidak atau".split())

def _tok(s: str) -> List[str]:
    return [w for w in re.findall(r"[a-z0-9]+", (s or "").lower()) if w not in STOP]

def _norm_choice(x: Optional[str]) -> str:
    if x is None:
        return ""
    s = str(x).strip().upper()
    for ch in s:
        if ch in CHOICES:
            return ch
    for ch in s:
        if ch in "12345":
            return "ABCDE"[int(ch) - 1]
    return ""

def _append_answers_csv(rows: List[dict]) -> None:
    path = JAWABAN_DIR / "siswa_jawaban.csv"
    path.parent.mkdir(parents=True, exist_ok=True)
    need_header = not path.exists() or path.stat().st_size == 0
    cols = ["student_id", "student_name", "session_id", "timestamp_ms", "mapel", "question_id", "chosen"]
    with path.open("a", newline="", encoding="utf-8") as fh:
        w = csv.DictWriter(fh, fieldnames=cols)
        if need_header:
            w.writeheader()
        for r in rows:
            w.writerow({k: r.get(k, "") for k in cols})

# ===== indexer & model loader =====
def _rebuild_indexes() -> None:
    """Bangun ulang index bank soal (dipanggil saat startup/reload)."""
    global BANKS, ANSKEY, ID2MAPEL
    BANKS = LC.load_all_banks()
    ANSKEY, ID2MAPEL = {}, {}
    for m, bank in BANKS.items():
        # Handle both list and dict formats
        questions = bank if isinstance(bank, list) else bank.get("soal", [])
        for q in questions:
            qid = str(q.get("id", ""))
            ANSKEY[(m, qid)] = {
                "kunci": (q.get("kunci") or "").strip().upper(),
                "bobot": int(q.get("bobot", 1)),
                "topik": q.get("topik") or "",
                "tingkat": q.get("tingkat") or "",
            }
            ID2MAPEL.setdefault(qid, set()).add(m)

def load_models() -> None:
    """Load semua model *.joblib di folder models/ ke memori."""
    global MODELS
    MODELS = {}
    if not MODELS_DIR.exists():
        return
    for p in MODELS_DIR.glob("clf_*.joblib"):
        mapel = p.stem.replace("clf_", "")
        try:
            MODELS[mapel.lower()] = joblib.load(p)
            print(f"[MODEL] loaded {mapel} <- {p.name}")
        except Exception as e:
            print(f"[MODEL] skip {p.name}: {e}")

def predict_ai(mapel: str, text: str, options: Dict[str, str]) -> Optional[str]:
    """Skor tiap opsi dengan model mapel; return huruf terbaik atau None jika model tidak ada."""
    mdl = MODELS.get((mapel or "").lower())
    if not mdl:
        return None
    cands, keys = [], []
    for k in ["A", "B", "C", "D", "E"]:
        if k in (options or {}):
            cands.append(f"{text} [SEP] {str(options.get(k, ''))}")
            keys.append(k)
    if not cands:
        return None
    # sklearn predict_proba -> list of [p(0), p(1)] → ambil p(1)
    probs = [p[1] for p in mdl.predict_proba(cands)]
    best_idx = max(range(len(probs)), key=lambda i: probs[i])
    return keys[best_idx] if keys else None

# ===== Material Processing Functions =====
def extract_text_from_pdf(file_path: Path) -> str:
    """Ekstrak teks dari file PDF."""
    try:
        import PyPDF2
        text = ""
        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            for page in pdf_reader.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"
        return text.strip()
    except Exception as e:
        print(f"[ERROR] Failed to extract PDF: {e}")
        import traceback
        traceback.print_exc()
        return ""

def extract_text_from_ppt(file_path: Path) -> str:
    """Ekstrak teks dari file PowerPoint."""
    try:
        from pptx import Presentation
        text = ""
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text.strip()
    except Exception as e:
        print(f"[ERROR] Failed to extract PPT: {e}")
        import traceback
        traceback.print_exc()
        return ""

def process_material_with_ai(material_text: str, mapel: str) -> dict:
    """Gunakan Gemini AI untuk membaca materi dan generate soal bank."""
    try:
        if not GOOGLE_GEMINI_API_KEY:
            print("[ERROR] Gemini API key not configured")
            return {"success": False, "error": "Gemini API key not configured", "questions": []}
        
        # Jalankan AI meski teks sedikit, dengan fallback chunking
        if not material_text or len(material_text.strip()) < 10:
            print(f"[ERROR] Material text too short: {len(material_text)} chars")
            return {"success": False, "error": "Material text is empty or too short", "questions": []}
        
        print(f"[AI] Processing material ({len(material_text)} chars) for {mapel}")
        
        # Ambil max 8000 chars; jika lebih, potong per 3000 chars dan gabungkan hasil
        chunks = []
        t = material_text.strip()
        MAX_INPUT = 8000
        if len(t) <= MAX_INPUT:
            chunks = [t]
        else:
            step = 3000
            for i in range(0, min(len(t), MAX_INPUT), step):
                chunks.append(t[i:i+step])

        prompt_tpl = f"""
Peran: Anda adalah guru ahli pembuat soal ujian untuk mata pelajaran {mapel}.

Tugas: Baca dan pelajari materi berikut, lalu buatlah 10-15 soal pilihan ganda berkualitas yang mencakup berbagai tingkat kesulitan (mudah, sedang, sulit) dan berbagai topik dalam materi.

Materi:
{{material_chunk}}

Instruksi:
1. Buat soal yang menguji pemahaman konsep, bukan hanya hafalan
2. Pastikan setiap soal memiliki 4-5 pilihan jawaban
3. Tandai jawaban yang benar
4. Berikan penjelasan singkat untuk setiap soal

Output HARUS dalam format JSON array (tanpa markdown, pure JSON):
[
  {{
    "question": "Teks pertanyaan...",
    "options": {{"A": "Pilihan A", "B": "Pilihan B", "C": "Pilihan C", "D": "Pilihan D", "E": "Pilihan E"}},
    "correct_answer": "A",
    "explanation": "Penjelasan mengapa A benar...",
    "difficulty": "mudah",
    "topic": "Topik spesifik dari materi"
  }}
]
"""
        
        print("[AI] Calling Gemini API...")
        if not CLIENT:
            raise RuntimeError("Gemini client not initialized")
        
        all_questions = []
        import time
        
        for idx, ch in enumerate(chunks):
            prompt = prompt_tpl.replace("{material_chunk}", ch)
            
            # Simple retry logic for 429 errors
            max_retries = 5
            retry_delay = 5 # seconds
            response = None
            
            for attempt in range(max_retries):
                try:
                    model_name = _pick_available_model()
                    response = CLIENT.models.generate_content(
                        model=model_name,
                        contents=prompt
                    )
                    break # Success, exit retry loop
                except Exception as e:
                    if "429" in str(e) or "RESOURCE_EXHAUSTED" in str(e):
                        if attempt < max_retries - 1:
                            print(f"[AI] Rate limit hit. Retrying in {retry_delay}s... (Attempt {attempt+1}/{max_retries})")
                            time.sleep(retry_delay)
                            retry_delay *= 2 # Exponential backoff
                        else:
                            print(f"[AI] Max retries exceeded for chunk {idx+1}")
                            raise e
                    else:
                        raise e # Re-raise other errors immediately

            if not response:
                continue

            response_text = (response.text or "").strip()
            print(f"[AI] Chunk {idx+1}/{len(chunks)} len={len(response_text)}")
            if not response_text:
                continue
            if response_text.startswith("```"):
                lines = response_text.split("\n")
                response_text = "\n".join(lines[1:-1])
                if response_text.startswith("json"):
                    response_text = response_text[4:].strip()
            try:
                data = json.loads(response_text)
                if isinstance(data, list):
                    all_questions.extend(data)
                elif isinstance(data, dict):
                    all_questions.append(data)
            except Exception as e:
                print(f"[AI] JSON parse failed on chunk {idx+1}: {e}")
                continue

        if not all_questions:
            return {"success": False, "error": "AI returned no questions", "questions": []}

        print(f"[AI] Successfully aggregated {len(all_questions)} questions")
        return {"success": True, "questions": all_questions[:15], "count": len(all_questions)}
        
    except json.JSONDecodeError as e:
        print(f"[ERROR] JSON parsing failed: {e}")
        return {"success": False, "error": f"JSON parse error: {str(e)}", "questions": []}
    except Exception as e:
        print(f"[ERROR] AI processing failed: {e}")
        import traceback
        traceback.print_exc()
        return {"success": False, "error": str(e), "questions": []}

def save_questions_to_bank(questions: list, mapel: str) -> int:
    """Save generated questions to JSON file AND PostgreSQL Database."""
    saved_count = 0
    
    # 1. Save to JSON File (Legacy/Backup)
    try:
        mapel_normalized = mapel.lower().replace(" ", "_")
        bank_file = SOAL_DIR / f"{mapel_normalized}.json"
        
        # Ensure SOAL_DIR exists
        SOAL_DIR.mkdir(parents=True, exist_ok=True)
        
        # Load existing bank or create new
        if bank_file.exists():
            with open(bank_file, 'r', encoding='utf-8') as f:
                bank_data = json.load(f)
        else:
            bank_data = []
        
        # Get next ID for JSON
        existing_ids = [int(q.get("id", 0)) for q in bank_data if isinstance(q.get("id"), (int, str)) and str(q.get("id")).isdigit()]
        next_id = max(existing_ids, default=0) + 1
        
        # Add new questions to JSON
        json_added = 0
        for q in questions:
            new_question = {
                "id": next_id,
                "teks": q.get("question", "") or q.get("pertanyaan", "") or q.get("teks", ""),
                "opsi": q.get("options", {}) or q.get("opsi", {}),
                "kunci": q.get("correct_answer", "") or q.get("jawaban_benar", "") or q.get("kunci", "A"),
                "topik": q.get("topic", "") or q.get("topik", ""),
                "tingkat": q.get("difficulty", "") or q.get("tingkat", "sedang"),
                "penjelasan": q.get("explanation", "") or q.get("pembahasan", ""),
                "source": "ai_material"
            }
            bank_data.append(new_question)
            next_id += 1
            json_added += 1
        
        # Save back to file
        with open(bank_file, 'w', encoding='utf-8') as f:
            json.dump(bank_data, f, ensure_ascii=False, indent=2)
        
        # Reload banks
        _rebuild_indexes()
        saved_count = json_added
        
    except Exception as e:
        print(f"[ERROR] Failed to save questions to JSON bank: {e}")

    # 2. Save to PostgreSQL Database (Persistence)
    try:
        db = next(get_db())
        db_added = 0
        
        for q in questions:
            # Normalize fields
            q_text = q.get("question", "") or q.get("pertanyaan", "") or q.get("teks", "")
            q_opsi = q.get("options", {}) or q.get("opsi", {})
            q_kunci = q.get("correct_answer", "") or q.get("jawaban_benar", "") or q.get("kunci", "A")
            
            if not q_text:
                continue

            # Check duplication
            existing = db.query(models.DBQuestion).filter(
                models.DBQuestion.mapel == mapel,
                models.DBQuestion.question_text == q_text
            ).first()
            
            if not existing:
                new_q = models.DBQuestion(
                    id=str(uuid.uuid4()),
                    mapel=mapel,
                    topic=q.get("topic", "") or q.get("topik", "Umum"),
                    difficulty=q.get("difficulty", "") or q.get("tingkat", "Sedang"),
                    question_text=q_text,
                    option_a=q_opsi.get("A", ""),
                    option_b=q_opsi.get("B", ""),
                    option_c=q_opsi.get("C", ""),
                    option_d=q_opsi.get("D", ""),
                    option_e=q_opsi.get("E", ""),
                    correct_answer=q_kunci,
                    explanation=q.get("explanation", "") or q.get("pembahasan", "") or q.get("penjelasan", ""),
                    created_at=datetime.utcnow()
                )
                db.add(new_q)
                db_added += 1
        
        db.commit()
        print(f"[DB] Successfully saved {db_added} questions to PostgreSQL.")
        
    except Exception as e:
        print(f"[ERROR] Failed to save questions to Database: {e}")
        
    return saved_count

@app.on_event("startup")
def _init():
    print("[STARTUP] Starting initialization...")
    
    try:
        print("[STARTUP] Loading indexes...")
        _rebuild_indexes()
        print(f"[STARTUP] Loaded {len(BANKS)} banks")
    except Exception as e:
        print(f"[STARTUP] Warning: Failed to load banks: {e}")
        # Continue anyway - quiz-settings will use defaults
    
    try:
        print("[STARTUP] Loading models...")
        load_models()
        print(f"[STARTUP] Loaded {len(MODELS)} models")
    except Exception as e:
        print(f"[STARTUP] Warning: Failed to load models: {e}")
    
    print("[STARTUP] Initialization complete")

# ===== models (Pydantic) =====
# Auth models
class LoginRequest(BaseModel):
    email: EmailStr
    password: str

class LoginResponse(BaseModel):
    token: str
    email: str
    username: str
    role: str
    subject: Optional[str] = None
    class_id: Optional[str] = None

class CreateUserRequest(BaseModel):
    email: EmailStr
    username: str
    password: str
    role: str  # "admin", "teacher", "student"
    subject: Optional[str] = None  # untuk teacher: "ipa", "ips", "matematika", etc
    class_id: Optional[str] = None  # untuk student

class UserResponse(BaseModel):
    email: str
    username: str
    role: str
    subject: Optional[str] = None
    class_id: Optional[str] = None

class CreateClassRequest(BaseModel):
    name: str  # contoh: "Kelas 7A"
    subject: str  # contoh: "ipa"

class ClassResponse(BaseModel):
    class_id: str
    name: str
    subject: str
    teacher_email: str
    teacher_name: str
    students: List[dict]  # [{email, username}]

class AssignStudentRequest(BaseModel):
    student_email: str
    class_id: str

class PredictRequest(BaseModel):
    text: str
    options: Dict[str, str]
    mapel: Optional[str] = None
    topik: Optional[str] = None

class PredictResponse(BaseModel):
    choice: str

class GenerateReq(BaseModel):
    mapel: Optional[List[str]] = None   # None → semua mapel
    n: int = Field(..., gt=0, description="Jumlah soal per mapel")
    topik: Optional[List[str]] = None
    tingkat: Optional[List[str]] = None
    seed: Optional[int] = None
    allow_duplicate: bool = False

class QuestionOut(BaseModel):
    mapel: str
    id: str
    topik: Optional[str] = ""
    tingkat: Optional[str] = ""
    teks: str
    opsi: Dict[str, str]  # A..E

class AnswerItem(BaseModel):
    mapel: Optional[str] = None
    question_id: str
    chosen: str

class ScoredItem(BaseModel):
    mapel: str
    question_id: str
    chosen: str
    kunci: str
    correct: bool
    bobot: int
    topik: Optional[str] = ""
    tingkat: Optional[str] = ""

class ScoreReq(BaseModel):
    student_id: str
    student_name: str
    session_id: Optional[str] = None
    answers: List[AnswerItem]
    materialize: bool = True  # buat artefak rekap guru (CSV) via evaluator

class ScoreResp(BaseModel):
    total: int
    benar: int
    persen: float
    per_mapel: Dict[str, Dict[str, float]]
    items: List[ScoredItem]

class RemediationTopic(BaseModel):
    topik: str
    mapel: str
    skor_persen: float
    is_remedial: bool

class RemediationRecommendation(BaseModel):
    student_id: str
    session_id: str
    total_score: float
    remedial_topics: List[RemediationTopic]
    recommended_actions: List[str]

def analyze_remedial(score_req: ScoreReq, mapel_scores: Dict, total_persen: float, session_id: str) -> RemediationRecommendation:
    """
    Analisis hasil ujian siswa.
    Tentukan topik mana yang perlu remedial (skor < 75%)
    """
    REMEDIAL_THERESHOLD = 0.75
    remedial_topics = []

    for mapel, stats in mapel_scores.items():
        persen = stats.get("persen", 0.0) / 100.0
        if persen < REMEDIAL_THERESHOLD:
            remedial_topics.append(RemediationTopic(
                topik=mapel,
                mapel=mapel,
                skor_persen=stats.get("persen", 0.0),
                is_remedial=True
            ))

    recommended_actions = []
    for rt in remedial_topics:
        recommended_actions.append(
            f"pelajari ulang materi {rt.mapel} - Skor: {rt.skor_persen:.1f}%"
        )

    return RemediationRecommendation(
        student_id=score_req.student_id,
        session_id=session_id,
        total_score=total_persen,
        remedial_topics=remedial_topics,
        recommended_actions=recommended_actions
    )

# ===== routes =====
@app.get("/")
def root():
    return {"service": "LearnCheck API", "docs": "/docs", "health": "/health"}

@app.get("/health")
def health():
    return {"ok": True, "banks": len(BANKS), "models": sorted(MODELS.keys())}

# ===== Database Ping =====
@app.get("/db/ping")
def db_ping(db: DBSessionType = Depends(get_db)):
    try:
        db.execute(text("SELECT 1"))
        return {"ok": True}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

# ===== Authentication & User Management =====
def get_current_user(credentials: HTTPAuthorizationCredentials = Depends(security), 
                     db: DBSessionType = Depends(get_db)) -> DBUser:
    """Verify token and return current user."""
    token = credentials.credentials
    session = crud.get_session(db, token)
    if not session:
        raise HTTPException(status_code=401, detail="Invalid or expired token")
    
    user = crud.get_user_by_email(db, session.email)
    if not user:
        raise HTTPException(status_code=401, detail="User not found")
    return user

def get_current_admin(credentials: HTTPAuthorizationCredentials = Depends(security),
                     db: DBSessionType = Depends(get_db)) -> DBUser:
    """Verify token and check admin role."""
    user = get_current_user(credentials, db)
    if user.role != "admin":
        raise HTTPException(status_code=403, detail="Admin access required")
    return user

@app.post("/auth/login", response_model=LoginResponse)
def login(req: LoginRequest, db: DBSessionType = Depends(get_db)):
    """Login endpoint for all users (admin, teacher, student)."""
    user = crud.get_user_by_email(db, req.email)
    if not user or not verify_password(req.password, user.hashed_password):
        raise HTTPException(status_code=401, detail="Invalid email or password")
    
    # Create session token
    token = create_token()
    crud.create_session(db, token, user.email, user.role)
    
    return LoginResponse(
        token=token,
        email=user.email,
        username=user.username,
        role=user.role,
        subject=user.subject,
        class_id=user.class_id
    )

@app.post("/auth/logout")
def logout(credentials: HTTPAuthorizationCredentials = Depends(security),
          db: DBSessionType = Depends(get_db)):
    """Logout - invalidate token."""
    token = credentials.credentials
    crud.delete_session(db, token)
    return {"message": "Logged out successfully"}

@app.post("/admin/users", response_model=UserResponse)
def create_user(req: CreateUserRequest, admin: DBUser = Depends(get_current_admin),
                db: DBSessionType = Depends(get_db)):
    """Admin endpoint to create new user (teacher or student)."""
    if crud.get_user_by_email(db, req.email):
        raise HTTPException(status_code=400, detail="Email already exists")
    
    if req.role not in ["admin", "teacher", "student"]:
        raise HTTPException(status_code=400, detail="Invalid role")
    
    # Validasi subject untuk teacher
    if req.role == "teacher" and not req.subject:
        raise HTTPException(status_code=400, detail="Teacher must have a subject")
    
    user = crud.create_user(
        db=db,
        email=req.email,
        username=req.username,
        hashed_password=hash_password(req.password),
        role=req.role,
        subject=req.subject
    )
    
    return UserResponse(
        email=user.email,
        username=user.username,
        role=user.role,
        subject=user.subject,
        class_id=user.class_id
    )

@app.get("/admin/users", response_model=List[UserResponse])
def list_users(admin: DBUser = Depends(get_current_admin),
               db: DBSessionType = Depends(get_db)):
    """Admin endpoint to list all users."""
    users = crud.get_all_users(db)
    return [
        UserResponse(
            email=u.email, 
            username=u.username, 
            role=u.role,
            subject=u.subject,
            class_id=u.class_id
        )
        for u in users
    ]

@app.delete("/admin/users/{email}")
def delete_user_endpoint(email: str, admin: DBUser = Depends(get_current_admin),
                        db: DBSessionType = Depends(get_db)):
    """Admin endpoint to delete a user."""
    if not crud.get_user_by_email(db, email):
        raise HTTPException(status_code=404, detail="User not found")
    
    if email == admin.email:
        raise HTTPException(status_code=400, detail="Cannot delete yourself")
    
    crud.delete_user(db, email)
    
    return {"message": "User deleted successfully"}

# ===== Class Management (Teacher) =====
@app.post("/teacher/classes", response_model=ClassResponse)
def create_class(req: CreateClassRequest, credentials: HTTPAuthorizationCredentials = Depends(security),
                db: DBSessionType = Depends(get_db)):
    """Teacher creates a new class for their subject."""
    try:
        teacher = get_current_user(credentials, db)
        if teacher.role != "teacher":
            raise HTTPException(status_code=403, detail="Only teachers can create classes")
        # Jika teacher.subject kosong (data lama), izinkan menggunakan req.subject
        teacher_subject = (teacher.subject or "").lower()
        req_subject = (req.subject or "").lower()
        if not teacher_subject:
            teacher_subject = req_subject
        if req_subject != teacher_subject:
            raise HTTPException(status_code=400, detail=f"You can only create classes for {teacher_subject or 'your subject'}")
        class_id = f"{req_subject}_{req.name}_{secrets.token_hex(4)}"
        cls = crud.create_class(
            db=db,
            class_id=class_id,
            name=req.name,
            subject=req_subject,
            teacher_email=teacher.email,
            teacher_name=teacher.username
        )
        return ClassResponse(
            class_id=cls.class_id,
            name=cls.name,
            subject=cls.subject,
            teacher_email=cls.teacher_email,
            teacher_name=cls.teacher_name,
            students=[]
        )
    except HTTPException:
        raise
    except Exception as e:
        print(f"[ERROR] create_class failed: {e}")
        raise HTTPException(status_code=500, detail="Internal error when creating class")

@app.get("/teacher/classes", response_model=List[ClassResponse])
def list_my_classes(credentials: HTTPAuthorizationCredentials = Depends(security),
                   db: DBSessionType = Depends(get_db)):
    """Teacher gets list of their classes."""
    teacher = get_current_user(credentials, db)
    
    if teacher.role != "teacher":
        raise HTTPException(status_code=403, detail="Only teachers can access this")
    
    classes = crud.get_classes_by_teacher(db, teacher.email)
    
    return [
        ClassResponse(
            class_id=c.class_id,
            name=c.name,
            subject=c.subject,
            teacher_email=c.teacher_email,
            teacher_name=c.teacher_name,
            students=[{"email": email, "username": email.split('@')[0]} 
                     for email in json.loads(c.students or "[]")]
        )
        for c in classes
    ]

@app.post("/teacher/classes/{class_id}/students")
def add_student_to_class(class_id: str, req: AssignStudentRequest, 
                        credentials: HTTPAuthorizationCredentials = Depends(security),
                        db: DBSessionType = Depends(get_db)):
    """Teacher adds a student to their class."""
    teacher = get_current_user(credentials, db)
    
    if teacher.role != "teacher":
        raise HTTPException(status_code=403, detail="Only teachers can add students")
    
    cls = crud.get_class_by_id(db, class_id)
    if not cls:
        raise HTTPException(status_code=404, detail="Class not found")
    
    # Validasi teacher owns this class
    if cls.teacher_email != teacher.email:
        raise HTTPException(status_code=403, detail="You don't own this class")
    
    # Validasi student exists
    student = crud.get_user_by_email(db, req.student_email)
    if not student:
        raise HTTPException(status_code=404, detail="Student not found")
    
    if student.role != "student":
        raise HTTPException(status_code=400, detail="User is not a student")
    
    # Add student to class
    crud.add_student_to_class(db, class_id, student.email)
    crud.update_user_class(db, student.email, class_id)
    
    # Reload class for response
    cls = crud.get_class_by_id(db, class_id)
    return {
        "message": "Student added successfully", 
        "class": ClassResponse(
            class_id=cls.class_id,
            name=cls.name,
            subject=cls.subject,
            teacher_email=cls.teacher_email,
            teacher_name=cls.teacher_name,
            students=[{"email": email, "username": email.split('@')[0]} 
                     for email in json.loads(cls.students or "[]")]
        )
    }

@app.delete("/teacher/classes/{class_id}/students/{student_email}")
def remove_student_from_class(class_id: str, student_email: str, 
                              credentials: HTTPAuthorizationCredentials = Depends(security),
                              db: DBSessionType = Depends(get_db)):
    """Teacher removes a student from their class."""
    teacher = get_current_user(credentials, db)
    
    if teacher.role != "teacher":
        raise HTTPException(status_code=403, detail="Only teachers can remove students")
    
    cls = crud.get_class_by_id(db, class_id)
    if not cls:
        raise HTTPException(status_code=404, detail="Class not found")
    
    if cls.teacher_email != teacher.email:
        raise HTTPException(status_code=403, detail="You don't own this class")
    
    # Remove student from class
    crud.remove_student_from_class(db, class_id, student_email)
    crud.update_user_class(db, student_email, None)
    
    return {"message": "Student removed successfully"}

@app.get("/teacher/available-students", response_model=List[UserResponse])
def get_available_students(credentials: HTTPAuthorizationCredentials = Depends(security),
                          db: DBSessionType = Depends(get_db)):
    """Teacher gets list of students not yet in any class."""
    teacher = get_current_user(credentials, db)
    
    if teacher.role != "teacher":
        raise HTTPException(status_code=403, detail="Only teachers can access this")
    
    # Get all students without class_id
    all_users = crud.get_all_users(db)
    available = [
        UserResponse(
            email=u.email,
            username=u.username,
            role=u.role,
            subject=u.subject,
            class_id=u.class_id
        )
        for u in all_users
        if u.role == "student" and not u.class_id
    ]
    
    return available

@app.post("/reload")
def reload_all():
    _rebuild_indexes()
    return {"ok": True, "banks": len(BANKS)}

@app.post("/models/reload")
def models_reload():
    load_models()
    return {"ok": True, "loaded": sorted(MODELS.keys())}

# ---- /predict: pakai AI jika ada model, else heuristik ----
@app.post("/predict", response_model=PredictResponse)
def predict(req: PredictRequest):
    if req.mapel:  # coba AI model spesifik mapel
        ai = predict_ai(req.mapel, req.text, req.options or {})
        if ai:
            return {"choice": ai}

    # fallback heuristik overlap token
    q = set(_tok(req.text))
    avail = [k for k in ["A", "B", "C", "D", "E"] if k in (req.options or {})] or ["A"]
    scores = {k: len(q & set(_tok(req.options.get(k, "")))) for k in avail}
    best = [k for k, v in scores.items() if v == max(scores.values())] or avail
    return {"choice": random.choice(best)}

# ---- generate quiz ----
@app.post("/generate_quiz", response_model=List[QuestionOut])
def generate_quiz(req: GenerateReq):
    targets = (sorted([p.stem for p in SOAL_DIR.glob("*.json")])
               if not req.mapel else [m.lower() for m in req.mapel])
    out: List[QuestionOut] = []
    for m in targets:
        bank = LC.load_bank_soal(m)
        cand = LC.filter_questions(bank, topik=req.topik, tingkat=req.tingkat)
        take = LC.pick_questions(cand, n=req.n, seed=req.seed, allow_duplicate=req.allow_duplicate)
        for q in take:
            opsi = q.get("opsi", {})
            out.append(QuestionOut(
                mapel=m,
                id=str(q.get("id", "")),
                topik=q.get("topik", ""),
                tingkat=q.get("tingkat", ""),
                teks=q.get("teks", ""),
                opsi={k: opsi.get(k, "") for k in ("A", "B", "C", "D", "E")},
            ))
    rng = random.Random(req.seed)
    rng.shuffle(out)
    return out

# ---- QG Service (Frontend Integration) ----
from training.scripts import qg_service

class QGRequest(BaseModel):
    question_text: str
    mapel: str
    topic: Optional[str] = ""
    difficulty: Optional[str] = "sedang"
    explanation: Optional[str] = ""
    choices: List[str] = []

@app.post("/qg/generate")
def qg_generate(req: QGRequest):
    """Endpoint untuk generate satu soal via Gemini (dipanggil dari Frontend)."""
    # Construct context
    context = f"Materi: {req.question_text}\nMapel: {req.mapel}\nTopik: {req.topic}\nKesulitan: {req.difficulty}"
    
    print(f"[QG] Generating question for {req.mapel}...")
    
    # Call Gemini via service
    raw_json = qg_service.generate_question_raw(context)
    
    # Clean markdown if present
    raw_json = raw_json.strip()
    if raw_json.startswith("```"):
        lines = raw_json.split("\n")
        # Remove first line
        lines = lines[1:]
        # Remove last line if it is just ```
        if lines and lines[-1].strip() == "```":
            lines = lines[:-1]
        raw_json = "\n".join(lines).strip()
    
    # Parse JSON
    try:
        data = json.loads(raw_json)
        
        # Frontend expects a single object, but service might return a list
        if isinstance(data, list) and len(data) > 0:
            return data[0] # Return first question
        elif isinstance(data, dict):
            return data
        else:
            # Fallback if empty list or other type
            raise ValueError("Empty or invalid format from AI")
            
    except Exception as e:
        print(f"[ERROR] Failed to parse QG response: {e}")
        # Fallback dummy so frontend doesn't break
        return {
            "question": "Gagal memproses respon AI (JSON Error). Silakan coba lagi.",
            "options": ["Opsi A", "Opsi B", "Opsi C", "Opsi D"],
            "answer_index": 0,
            "explanation": f"Error detail: {str(e)}"
        }

# ---- simple generate endpoint (for frontend) ----
@app.post("/generate")
def generate_simple(mapel: str = None, n: int = 10):
    """Simple endpoint to generate quiz questions from database."""
    if not mapel:
        return {"error": "Parameter 'mapel' diperlukan"}
    
    mapel_lower = mapel.lower()
    bank = BANKS.get(mapel_lower)
    
    # Try finding case-insensitive match if exact match fails
    if not bank:
        for k in BANKS.keys():
            if k.lower() == mapel_lower:
                bank = BANKS[k]
                break
    
    if not bank:
        # Debug info
        print(f"[GENERATE] Mapel '{mapel}' not found. Available: {list(BANKS.keys())}")
        # Try to reload banks just in case
        _rebuild_indexes()
        bank = BANKS.get(mapel_lower)
        
        if not bank:
             return {"error": f"Mapel '{mapel}' tidak ditemukan. Available: {list(BANKS.keys())}"}
    
    # Pick random questions (no allow_duplicate parameter)
    questions = LC.pick_questions(bank, n=n)
    
    out = []
    for q in questions:
        opsi = q.get("opsi", {})
        # Map field names: "pertanyaan" or "teks" -> "teks"
        question_text = q.get("pertanyaan") or q.get("teks", "")
        
        # Ensure we have a valid question text
        if not question_text:
            continue
            
        # Ensure options exist
        if not opsi:
            continue

        out.append({
            "id": str(q.get("id", "")),
            "mapel": mapel,
            "teks": question_text,
            "opsi": {k: opsi.get(k, "") for k in ("A", "B", "C", "D", "E")},
            "topik": q.get("topik", ""),
            "tingkat": q.get("tingkat", ""),
            "jawaban_benar": q.get("kunci") or q.get("jawaban_benar", ""),  # Include correct answer key
        })
    
    if not out:
        return {"error": "Tidak ada soal yang valid ditemukan untuk mapel ini"}

    return out

# ---- score (tanpa remedial) + rekap guru opsional ----
@app.post("/score", response_model=ScoreResp)
def score(req: ScoreReq):
    session_id = req.session_id or f"{req.student_id}-{int(datetime.utcnow().timestamp() * 1000)}"
    items: List[ScoredItem] = []
    per_mapel: Dict[str, Dict[str, float]] = {}
    to_log = []
    benar = 0

    for a in req.answers:
        qid = str(a.question_id).strip()
        m = (a.mapel or next(iter(ID2MAPEL.get(qid, [])), "")).lower()
        ak = ANSKEY.get((m, qid))
        chosen = _norm_choice(a.chosen)

        if not ak:
            items.append(ScoredItem(
                mapel=m, question_id=qid, chosen=chosen, kunci="",
                correct=False, bobot=0, topik="", tingkat=""
            ))
            continue

        kunci = ak["kunci"]
        bobot = int(ak["bobot"])
        ok = (chosen == kunci) and bool(kunci)
        if ok:
            benar += 1

        items.append(ScoredItem(
            mapel=m, question_id=qid, chosen=chosen, kunci=kunci,
            correct=ok, bobot=bobot, topik=ak["topik"], tingkat=ak["tingkat"]
        ))

        st = per_mapel.setdefault(m, {"benar": 0, "total": 0, "bobot": 0.0, "skor": 0.0})
        st["total"] += 1
        st["bobot"] += bobot
        if ok:
            st["benar"] += 1
            st["skor"] += bobot

        to_log.append({
            "student_id": req.student_id,
            "student_name": req.student_name,
            "session_id": session_id,
            "timestamp_ms": str(int(datetime.utcnow().timestamp() * 1000)),
            "mapel": m,
            "question_id": qid,
            "chosen": chosen,
        })

    # simpan jawaban siswa
    _append_answers_csv(to_log)

    # (opsional) buat artefak rekap guru dengan evaluator yang sudah ada
    if req.materialize:
        out_dir = JAWABAN_DIR / "out"
        try:
            subprocess.run(
                [sys.executable, "-m", "training.scripts.evaluate_answers",
                 "--in", str(JAWABAN_DIR / "siswa_jawaban.csv"),
                 "--out-dir", str(out_dir)],
                check=True
            )
        except Exception as e:
            print(f"[WARN] evaluate_answers gagal: {e}")

    total = len(items)
    persen = round(100.0 * benar / total, 2) if total else 0.0
    per_mapel_pct = {
        m: {"benar": v["benar"], "total": v["total"],
            "persen": (round(100.0 * v["skor"] / v["bobot"], 2) if v["bobot"] else 0.0)}
        for m, v in per_mapel.items()
    }

    remedial = analyze_remedial(req, per_mapel_pct, persen, session_id)

    # TODO: implement _save_remedial_recommendation
    # if remedial.remedial_topics:
    #     _save_remedial_recommendation(remedial)

    return ScoreResp(total=total, benar=benar, persen=persen,
                     per_mapel=per_mapel_pct, items=items)

@app.get("/remedial/{student_id}/{session_id}")
def generate_remedial_quiz(
    mapel: str,
    n: int = 5,
    tingkat: str = "mudah" # Remedial = mudah
):
    """ Generate soal level MUDAH untuk remedial."""
    bank = LC.load_bank_soal(mapel)
    cand = LC.filter_questions(bank, tingkat=[tingkat])
    take = LC.pick_questions(cand, n=n)

    out = []
    for q in take:
        opsi = q.get("opsi", {})
        out.append(QuestionOut(
            mapel=mapel,
            id=str(q.get("id", "")),
            topik=q.get("tingkat", ""),
            teks = q.get("teks", ""),
            opsi={k: opsi.get(k, "") for k in ("A", "B", "C", "D", "E")},
        ))
    return out

# Remedial recommendation endpoint
class RemedialRequest(BaseModel):
    mapel: str
    wrong_questions: List[str]

@app.post("/remedial/recommend")
def recommend_remedial(req: RemedialRequest):
    """Generate remedial content recommendations based on wrong answers."""
    try:
        # Load materi file
        mapel_file = req.mapel.lower().replace(" ", "_")
        materi_path = Path("training/materi") / f"{mapel_file}.txt"
        
        if not materi_path.exists():
            return {"content": f"Materi untuk {req.mapel} belum tersedia. Silakan hubungi guru untuk materi pembelajaran.", "mapel": req.mapel}
        
        # Read materi content
        with open(materi_path, 'r', encoding='utf-8') as f:
            materi_content = f.read()
        
        # Gunakan Gemini untuk generate rekomendasi remedial
        if not GOOGLE_GEMINI_API_KEY:
            print("[ERROR] Gemini API key not configured")
            return {"content": "Gemini API key not configured.", "mapel": req.mapel}

        prompt = f"""Anda adalah asisten pembelajaran yang membantu siswa memahami materi {req.mapel}.

    Siswa telah mengerjakan quiz dan salah menjawab pertanyaan berikut:
    {chr(10).join(f"- {q}" for q in req.wrong_questions)}

    Berikut adalah materi lengkap untuk {req.mapel}:
    {materi_content[:8000]}

    Tugas Anda:
    1. Analisis pertanyaan yang salah dan identifikasi topik/konsep yang perlu dipelajari ulang
    2. Berikan rekomendasi materi SPESIFIK dari materi yang tersedia yang relevan dengan kesalahan siswa
    3. Kutip bagian materi yang paling penting untuk dipelajari ulang
    4. Berikan tips praktis untuk memahami konsep tersebut

    Format output (gunakan Markdown):
    ## 📚 Rekomendasi Materi Remedial

    ### Topik yang Perlu Dipelajari Ulang:
    [List topik dari analisis kesalahan]

    ### Materi yang Relevan:
    [Kutip materi spesifik dari konten yang diberikan, fokus pada konsep yang belum dipahami]

    ### Tips Belajar:
    [Berikan 2-3 tips praktis untuk memahami konsep]

    Berikan penjelasan yang jelas, spesifik, dan mudah dipahami siswa."""

        print("[AI] Calling Gemini API for remedial...")
        if not CLIENT:
            raise RuntimeError("Gemini client not initialized")
        response = CLIENT.models.generate_content(
            model="gemini-1.5-flash",
            contents=prompt
        )
        content = response.text
        print(f"[AI] Remedial response: {content[:200]}")
        return {"content": content, "mapel": req.mapel}
        
    except Exception as e:
        print(f"Error in remedial recommend: {e}")
        return {"content": f"Materi remedial: Silakan review kembali topik {req.mapel.replace('_', ' ').title()}. Fokus pada konsep dasar dan latihan soal tambahan.", "error": str(e)}

# ===== Material Management Endpoints =====
@app.post("/materials/upload")
async def upload_material(
    file: UploadFile = File(...),
    title: str = Form(""),
    description: str = Form(""),
    subject: str = Form(""),
    db: DBSessionType = Depends(get_db)
):
    """Upload materi pembelajaran (PDF/PPT) dan proses dengan AI untuk generate soal."""
    try:
        print(f"\n[UPLOAD] ========== NEW UPLOAD ==========")
        print(f"[UPLOAD] File: {file.filename}")
        print(f"[UPLOAD] Subject: {subject}")
        print(f"[UPLOAD] Title: {title}")
        print(f"[UPLOAD] Content-Type: {file.content_type}")
        
        # Normalize subject (case insensitive, spaces to underscores)
        subject_normalized = subject.lower().replace(" ", "_")
        print(f"[UPLOAD] Subject normalized: {subject_normalized}")
        
        # Validate file type
        allowed_types = ["application/pdf", "application/vnd.ms-powerpoint", 
                        "application/vnd.openxmlformats-officedocument.presentationml.presentation"]
        
        if file.content_type not in allowed_types:
            print(f"[ERROR] Invalid file type: {file.content_type}")
            raise HTTPException(status_code=400, detail=f"Only PDF or PPT files are allowed. Got: {file.content_type}")
        
        # Validate file size (max 10MB)
        file_content = await file.read()
        print(f"[UPLOAD] File size: {len(file_content)} bytes ({len(file_content)/1024/1024:.2f} MB)")
        
        if len(file_content) > 10 * 1024 * 1024:
            raise HTTPException(status_code=400, detail="File size exceeds 10MB limit")
        
        # Save file to uploads directory
        file_extension = file.filename.split(".")[-1].lower()
        safe_filename = f"{subject_normalized}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.{file_extension}"
        file_path = UPLOADS_DIR / safe_filename
        
        print(f"[UPLOAD] Saving to: {file_path}")
        with open(file_path, "wb") as f:
            f.write(file_content)
        print(f"[UPLOAD] File saved successfully")
        

        # Ekstrak teks dari file (PDF/PPT)
        material_text = ""
        if file_extension == "pdf":
            material_text = extract_text_from_pdf(file_path)
            if not material_text or len(material_text.strip()) < 10:
                print("[UPLOAD] PDF text extraction failed or too short, trying OCR...")
                material_text = extract_text_from_pdf_with_ocr(file_path)
        elif file_extension in ["ppt", "pptx"]:
            material_text = extract_text_from_ppt(file_path)
        # Jika ingin support gambar: elif file_extension in ["jpg", "jpeg", "png"]: ...
        print(f"[UPLOAD] Extracted {len(material_text)} chars from materi (with OCR fallback if needed)")

        # Simpan ke database tanpa cek material_text
        
        # Save to database
        print(f"[UPLOAD] Saving to database...")
        try:
            # Get teacher/uploader ID from auth token (for now, use default admin)
            uploader = db.query(DBUser).filter(DBUser.role == "admin").first()
            if not uploader:
                print("[ERROR] No admin user found in database")
                raise HTTPException(status_code=500, detail="No admin user found")
            
            print(f"[UPLOAD] Using uploader: {uploader.email}")
            
            new_material = DBMaterial(
                title=title or file.filename,
                description=description or f"Materi {subject}",
                file_url=f"/uploads/{safe_filename}",
                file_type=file_extension,
                mapel=subject_normalized,
                uploader_id=uploader.id
            )
            
            db.add(new_material)
            db.commit()
            db.refresh(new_material)
            print(f"[UPLOAD] ✅ Material saved to DB with ID: {new_material.id}")
            
        except Exception as db_error:
            print(f"[ERROR] Database save failed: {db_error}")
            import traceback
            traceback.print_exc()
            db.rollback()
            raise HTTPException(status_code=500, detail=f"Database error: {str(db_error)}")
        
        # Proses AI untuk generate soal dari materi
        jumlah_soal = 0
        ai_result = None
        if material_text and len(material_text.strip()) >= 50:
            ai_result = process_material_with_ai(material_text, subject_normalized)
            if ai_result and ai_result.get("success") and ai_result.get("questions"):
                jumlah_soal = save_questions_to_bank(ai_result["questions"], subject_normalized)
            else:
                print(f"[UPLOAD] AI gagal generate soal: {ai_result.get('error') if ai_result else 'No result'}")
        else:
            print("[UPLOAD] Materi kosong atau terlalu pendek, AI tidak dijalankan.")

        print(f"[UPLOAD] ========== UPLOAD COMPLETE ==========")
        return {
            "success": True,
            "material_id": new_material.id,
            "title": new_material.title,
            "mapel": new_material.mapel,
            "file_path": str(file_path),
            "jumlah_soal_ditambahkan": jumlah_soal
        }
        
    except HTTPException:
        raise
    except Exception as e:
        print(f"[ERROR] Material upload failed: {e}")
        raise HTTPException(status_code=500, detail=f"Upload failed: {str(e)}")

@app.get("/materials")
def list_materials(
    mapel: Optional[str] = None,
    db: DBSessionType = Depends(get_db)
):
    """Get list of all materials, optionally filtered by mapel."""
    try:
        query = db.query(DBMaterial)
        
        if mapel:
            # Use ilike for case-insensitive search
            query = query.filter(DBMaterial.mapel.ilike(mapel))
        
        materials = query.order_by(DBMaterial.created_at.desc()).all()
        
        result = []
        for m in materials:
            result.append({
                "id": m.id,
                "title": m.title,
                "description": m.description,
                "file_url": m.file_url,
                "file_type": m.file_type,
                "mapel": m.mapel,
                "uploaded_by": m.uploader_id,
                "created_at": m.created_at.isoformat() if m.created_at else None
            })
        
        return result
        
    except Exception as e:
        print(f"[ERROR] Failed to list materials: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/materials/{material_id}")
def get_material(material_id: int, db: DBSessionType = Depends(get_db)):
    """Get specific material by ID."""
    material = db.query(DBMaterial).filter(DBMaterial.id == material_id).first()
    
    if not material:
        raise HTTPException(status_code=404, detail="Material not found")
    
    return {
        "id": material.id,
        "title": material.title,
        "description": material.description,
        "file_url": material.file_url,
        "file_type": material.file_type,
        "mapel": material.mapel,
        "uploaded_by": material.uploader_id,
        "created_at": material.created_at.isoformat() if material.created_at else None
    }

@app.delete("/materials/{material_id}")
def delete_material(material_id: int, db: DBSessionType = Depends(get_db)):
    material = db.query(DBMaterial).filter(DBMaterial.id == material_id).first()
    if not material:
        raise HTTPException(status_code=404, detail="Material not found")
    db.delete(material)
    db.commit()
    return {"success": True, "message": "Material deleted"}

# ===== Quiz Settings Management =====
class QuizSettingsReq(BaseModel):
    mapel: str
    enabled: bool = True
    timer: int = 60
    startDate: Optional[str] = None
    endDate: Optional[str] = None
    showCorrectAnswers: bool = True
    randomizeQuestions: bool = False
    attempts: int = 1

class QuizSettingsResp(BaseModel):
    mapel: str
    enabled: bool
    timer: int
    startDate: Optional[str]
    endDate: Optional[str]
    showCorrectAnswers: bool
    randomizeQuestions: bool
    attempts: int

@app.get("/quiz-settings/{mapel}")
def get_quiz_settings_endpoint(mapel: str, db: DBSessionType = Depends(get_db)) -> QuizSettingsResp:
    """Get quiz settings for a specific mapel."""
    quiz = crud.get_quiz_settings(db, mapel.lower())
    
    if quiz:
        return QuizSettingsResp(
            mapel=mapel,
            enabled=quiz.enabled,
            timer=quiz.duration,
            startDate=quiz.start_date.isoformat() if quiz.start_date else None,
            endDate=quiz.end_date.isoformat() if quiz.end_date else None,
            showCorrectAnswers=quiz.show_correct_answers,
            randomizeQuestions=quiz.randomize_questions,
            attempts=quiz.max_attempts
        )
    else:
        # Return defaults
        return QuizSettingsResp(
            mapel=mapel,
            enabled=True,
            timer=60,
            startDate=None,
            endDate=None,
            showCorrectAnswers=True,
            randomizeQuestions=False,
            attempts=1
        )

@app.post("/quiz-settings")
def update_quiz_settings_endpoint(req: QuizSettingsReq, db: DBSessionType = Depends(get_db)) -> QuizSettingsResp:
    """Update quiz settings for a specific mapel."""
    from datetime import datetime
    
    start_date = None
    end_date = None
    
    if req.startDate:
        try:
            start_date = datetime.fromisoformat(req.startDate.replace('Z', '+00:00'))
        except:
            pass
    
    if req.endDate:
        try:
            end_date = datetime.fromisoformat(req.endDate.replace('Z', '+00:00'))
        except:
            pass
    
    quiz = crud.upsert_quiz_settings(
        db=db,
        mapel=req.mapel.lower(),
        enabled=req.enabled,
        timer=req.timer,
        start_date=start_date,
        end_date=end_date,
        show_correct_answers=req.showCorrectAnswers,
        randomize_questions=req.randomizeQuestions,
        attempts=req.attempts
    )
    
    return QuizSettingsResp(
        mapel=req.mapel,
        enabled=quiz.enabled,
        timer=quiz.duration,
        startDate=quiz.start_date.isoformat() if quiz.start_date else None,
        endDate=quiz.end_date.isoformat() if quiz.end_date else None,
        showCorrectAnswers=quiz.show_correct_answers,
        randomizeQuestions=quiz.randomize_questions,
        attempts=quiz.max_attempts
    )

@app.get("/quiz-settings")
def get_all_quiz_settings_endpoint(db: DBSessionType = Depends(get_db)) -> List[QuizSettingsResp]:
    """Get quiz settings for all mapels."""
    # Default mapels
    default_mapels = [
        "agama_hindu", "agama_islam", "agama_kristen", 
        "biologi", "ekonomi", "fisika", "geografi", 
        "ipa", "ips", "kesenian", "kimia", "matematika",
        "penjaskes", "ppkn", "sejarah", "sosiologi"
    ]
    
    # Use BANKS if available, otherwise use defaults
    mapel_list = list(BANKS.keys()) if BANKS else default_mapels
    
    result = []
    for mapel in mapel_list:
        quiz = crud.get_quiz_settings(db, mapel.lower())
        
        if quiz:
            result.append(QuizSettingsResp(
                mapel=mapel,
                enabled=quiz.enabled,
                timer=quiz.duration,
                startDate=quiz.start_date.isoformat() if quiz.start_date else None,
                endDate=quiz.end_date.isoformat() if quiz.end_date else None,
                showCorrectAnswers=quiz.show_correct_answers,
                randomizeQuestions=quiz.randomize_questions,
                attempts=quiz.max_attempts
            ))
        else:
            # Return defaults for mapels without settings
            result.append(QuizSettingsResp(
                mapel=mapel,
                enabled=True,
                timer=60,
                startDate=None,
                endDate=None,
                showCorrectAnswers=True,
                randomizeQuestions=False,
                attempts=1
            ))
    
    return result


# ===== Run server =====
if __name__ == "__main__":
    import uvicorn
    print("[STARTUP] Starting server on http://localhost:8000")
    uvicorn.run("main:app", host="0.0.0.0", port=8000, reload=True)
